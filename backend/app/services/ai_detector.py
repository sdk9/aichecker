"""
AI-generated content detection service.

Detection layers (in order of reliability):
  1. ML Neural-Network classifier  — primary signal when available
  2. Forensic heuristics           — corroborating evidence

Supported file types:
  Images        — JPG / PNG / WebP / GIF / TIFF / BMP / HEIC
  Documents     — PDF / DOCX / DOC
  Presentations — PPTX / PPT
  Spreadsheets  — XLSX / XLS / CSV

Weights when ML is available:
  Images    — ML 60 % + heuristics 40 %
  Text docs — RoBERTa 15 % + GPT-2 perplexity 20 % + heuristics 65 %

Weights when ML is NOT available (torch not installed):
  Images    — ELA 28 % + EXIF 18 % + Noise 22 % + Freq 16 % + Color 9 % + Texture 7 %
  Text docs — Template 18 % + Buzzwords 14 % + CV structure 12 % + Achievements 10 % +
              Skills 9 % + Burstiness 10 % + Vocab 8 % + Phrases 10 % + Para 5 % + Sim 4 %
"""

import io
import logging
import math
import re
import zipfile
from pathlib import Path
from typing import List, Optional, Tuple
from xml.etree import ElementTree as ET

from app.models.analysis import DetectionSignal, SignalSeverity

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# IMAGE ANALYSIS
# ─────────────────────────────────────────────────────────────────────────────

async def analyze_image(file_path: Path) -> Tuple[float, List[DetectionSignal]]:
    """Returns (ai_probability 0-1, signals)."""
    signals: List[DetectionSignal] = []

    try:
        from PIL import Image
        import numpy as np

        img = Image.open(file_path).convert("RGB")
        arr = np.array(img, dtype=np.float32)

        # ── 1. ML classifier (primary signal) ─────────────────────────────
        ml_prob, ml_weight = _ml_image_signal(file_path, signals)

        # ── 2-7. Forensic heuristics ───────────────────────────────────────
        exif_score,  exif_sig  = _check_exif_presence(file_path, img)
        ela_score,   ela_sig   = _ela_analysis(img)
        noise_score, noise_sig = _noise_analysis(arr)
        freq_score,  freq_sig  = _frequency_analysis(arr)
        color_score, color_sig = _color_distribution(arr)
        tex_score,   tex_sig   = _texture_regularity(arr)

        signals += [exif_sig, ela_sig, noise_sig, freq_sig, color_sig, tex_sig]
        h_scores = [exif_score, ela_score, noise_score, freq_score, color_score, tex_score]

        if ml_prob is not None:
            # ML available: 60 % ML + 40 % heuristics
            h_weights = [0.12, 0.10, 0.08, 0.06, 0.02, 0.02]  # sum = 0.40
            h_prob = sum(s * w for s, w in zip(h_scores, h_weights))
            ai_prob = ml_prob * ml_weight + h_prob
        else:
            # Heuristics only
            h_weights = [0.18, 0.28, 0.22, 0.16, 0.09, 0.07]  # sum = 1.00
            ai_prob = sum(s * w for s, w in zip(h_scores, h_weights))

    except Exception as exc:
        logger.exception("Image analysis failed")
        signals.append(DetectionSignal(
            name="Image Analysis Error",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.0,
        ))
        return 0.5, signals

    return float(min(max(ai_prob, 0.0), 1.0)), signals


def _ml_image_signal(
    file_path: Path,
    signals: List[DetectionSignal],
) -> Tuple[Optional[float], float]:
    """
    Attempt ML classification. Appends a signal and returns
    (ml_probability, ml_weight) or (None, 0) if unavailable.
    """
    try:
        from app.services.ml_detector import ml_image_score
        ml_prob, details = ml_image_score(file_path)
        if ml_prob is None:
            return None, 0.0

        if ml_prob >= 0.70:
            severity = SignalSeverity.HIGH
        elif ml_prob >= 0.40:
            severity = SignalSeverity.MEDIUM
        else:
            severity = SignalSeverity.LOW

        signals.append(DetectionSignal(
            name="Neural Network Classifier",
            description=(
                f"Vision Transformer trained on hundreds of thousands of AI-generated "
                f"and real images. Confidence: {ml_prob:.1%} AI-generated."
            ),
            severity=severity,
            score=ml_prob,
            details=details,
        ))
        return ml_prob, 0.60

    except Exception as exc:
        logger.warning("ML image signal failed: %s", exc)
        return None, 0.0


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENT ANALYSIS
# ─────────────────────────────────────────────────────────────────────────────

async def analyze_document(file_path: Path, mime_type: str) -> Tuple[float, List[DetectionSignal]]:
    """Returns (ai_probability, signals) for PDF and Word documents."""
    text = ""
    try:
        if mime_type == "application/pdf":
            text = _extract_pdf_text(file_path)
            ai_prob, signals = _analyse_text_content(text)
            pdf_score, pdf_sig = _pdf_package_signal(file_path, text)
            signals.append(pdf_sig)
            if pdf_score >= 0.68:
                ai_prob = max(ai_prob, min(ai_prob * 0.85 + pdf_score * 0.15, 1.0))
            if pdf_score >= 0.90:
                ai_prob = max(ai_prob, 0.84)
            elif pdf_score >= 0.78:
                ai_prob = max(ai_prob, 0.72)
            elif pdf_score >= 0.68:
                ai_prob = max(ai_prob, 0.62)
            return ai_prob, signals
        else:
            text = _extract_docx_text(file_path)
            ai_prob, signals = _analyse_text_content(text)
            docx_score, docx_sig = _docx_package_signal(file_path, text)
            signals.append(docx_sig)
            # Only boost on suspicious metadata — never let clean metadata pull text score down
            if docx_score >= 0.72:
                ai_prob = max(ai_prob, min(ai_prob * 0.88 + docx_score * 0.12, 1.0))
            if docx_score >= 0.90:
                ai_prob = max(ai_prob, 0.88)
            elif docx_score >= 0.72:
                ai_prob = max(ai_prob, 0.74)
            return ai_prob, signals
    except Exception as exc:
        logger.warning("Document extraction failed: %s", exc)
    return _analyse_text_content(text)


def _ml_text_signal(
    text: str,
    signals: List[DetectionSignal],
) -> Tuple[Optional[float], float]:
    try:
        from app.services.ml_detector import ml_text_score
        ml_prob, details = ml_text_score(text)
        if ml_prob is None:
            return None, 0.0

        if ml_prob >= 0.70:
            severity = SignalSeverity.HIGH
        elif ml_prob >= 0.40:
            severity = SignalSeverity.MEDIUM
        else:
            severity = SignalSeverity.LOW

        signals.append(DetectionSignal(
            name="Neural Network Text Classifier",
            description=(
                f"RoBERTa model fine-tuned to distinguish human vs AI-generated text. "
                f"Confidence: {ml_prob:.1%} AI-generated."
            ),
            severity=severity,
            score=ml_prob,
            details=details,
        ))
        return ml_prob, 0.55

    except Exception as exc:
        logger.warning("ML text signal failed: %s", exc)
        return None, 0.0


# ─────────────────────────────────────────────────────────────────────────────
# PRESENTATION ANALYSIS (PPTX / PPT)
# ─────────────────────────────────────────────────────────────────────────────

async def analyze_presentation(file_path: Path, mime_type: str) -> Tuple[float, List[DetectionSignal]]:
    """Returns (ai_probability, signals) for PowerPoint presentations."""
    text = ""
    slide_texts: List[str] = []
    try:
        text, slide_texts = _extract_pptx_text(file_path)
    except Exception as exc:
        logger.warning("PPTX extraction failed: %s", exc)
    ai_prob, signals = _analyse_text_content(text)
    if len(slide_texts) >= 3:
        slide_score, slide_sig = _slide_structural_uniformity(slide_texts)
        signals.append(slide_sig)
        ai_prob = min(float(ai_prob) * 0.88 + float(slide_score) * 0.12, 1.0)
    return float(min(max(ai_prob, 0.0), 1.0)), signals


# ─────────────────────────────────────────────────────────────────────────────
# SPREADSHEET ANALYSIS (XLSX / XLS / CSV)
# ─────────────────────────────────────────────────────────────────────────────

async def analyze_spreadsheet(file_path: Path, mime_type: str) -> Tuple[float, List[DetectionSignal]]:
    """Returns (ai_probability, signals) for Excel and CSV files."""
    text = ""
    try:
        if mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            text = _extract_xlsx_text(file_path)
        else:
            text = _extract_csv_text(file_path)
    except Exception as exc:
        logger.warning("Spreadsheet extraction failed: %s", exc)
    return _analyse_text_content(text)


# ─────────────────────────────────────────────────────────────────────────────
# IMAGE FORENSIC HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _check_exif_presence(file_path: Path, img) -> Tuple[float, DetectionSignal]:
    """
    AI generators almost never embed authentic camera EXIF.

    Calibration note:
      - "No EXIF" is a *weak* signal (screenshots, social-media images, web downloads
        all lack EXIF legitimately). Cap at 0.60 to reduce false positives.
      - "Full camera EXIF" is a strong negative signal (very hard to fake convincingly).
    """
    try:
        import exifread
        with open(file_path, "rb") as f:
            tags = exifread.process_file(f, details=False)

        has_camera = any("Make" in k or "Model" in k for k in tags)
        has_gps    = any("GPS" in k for k in tags)
        has_lens   = any("Lens" in k for k in tags)
        has_datetime = any("DateTime" in k for k in tags)

        if not tags:
            return 0.60, DetectionSignal(
                name="No EXIF Metadata",
                description=(
                    "File has no EXIF data. Authentic camera photos always contain "
                    "hardware metadata; AI generators do not embed EXIF."
                ),
                severity=SignalSeverity.MEDIUM,
                score=0.60,
                details="Note: screenshots and socially-shared images also lack EXIF.",
            )

        if not has_camera:
            return 0.50, DetectionSignal(
                name="Missing Camera Hardware Info",
                description="EXIF present but lacks camera make/model — unusual for real photos.",
                severity=SignalSeverity.MEDIUM,
                score=0.50,
            )

        detail = f"Camera: {tags.get('Image Make', '?')} {tags.get('Image Model', '?')}"
        if has_gps:
            detail += " | GPS data present"
        if has_lens:
            detail += " | Lens info present"
        if has_datetime:
            detail += f" | Taken: {tags.get('Image DateTime', '?')}"

        return 0.08, DetectionSignal(
            name="Camera EXIF Present",
            description="Full camera hardware metadata found. Consistent with authentic photography.",
            severity=SignalSeverity.LOW,
            score=0.08,
            details=detail,
        )
    except Exception:
        return 0.45, DetectionSignal(
            name="EXIF Check Failed",
            description="Could not parse EXIF data.",
            severity=SignalSeverity.LOW,
            score=0.45,
        )


def _ela_analysis(img) -> Tuple[float, DetectionSignal]:
    """
    Error Level Analysis — resave at known JPEG quality and measure pixel delta.
    AI-generated images have characteristically uniform ELA across the whole image.

    Fixed vs previous version:
      - PNG/non-JPEG inputs: the first JPEG save introduces pervasive compression
        artifacts, so we use a higher quality (92) to reduce that artefact and
        recalibrate thresholds accordingly.
    """
    try:
        import numpy as np
        from PIL import Image

        img_rgb = img.convert("RGB")
        is_jpeg = getattr(img, "format", "") in ("JPEG", "JPG")
        quality = 75 if is_jpeg else 92  # higher quality for non-JPEG to reduce noise

        buf = io.BytesIO()
        img_rgb.save(buf, format="JPEG", quality=quality)
        buf.seek(0)
        img_resaved = Image.open(buf).convert("RGB")

        orig = np.array(img_rgb,  dtype=np.float32)
        resd = np.array(img_resaved, dtype=np.float32)
        ela  = np.abs(orig - resd)

        mean_ela = float(ela.mean())
        std_ela  = float(ela.std())

        h, w = ela.shape[:2]
        block = 32
        block_means = [
            ela[y:y + block, x:x + block].mean()
            for y in range(0, h - block, block)
            for x in range(0, w - block, block)
        ]
        block_cv = float(
            np.array(block_means).std() / (np.array(block_means).mean() + 1e-6)
        )

        # Recalibrated thresholds (less aggressive than v1)
        if block_cv < 0.35 and mean_ela < 10.0:
            score, severity = 0.72, SignalSeverity.HIGH
            desc = (
                f"Suspiciously uniform ELA (mean={mean_ela:.1f}, CV={block_cv:.2f}). "
                "AI images lack natural JPEG re-compression variance."
            )
        elif block_cv < 0.60 or mean_ela < 6.0:
            score, severity = 0.42, SignalSeverity.MEDIUM
            desc = (
                f"Moderately uniform ELA (mean={mean_ela:.1f}, CV={block_cv:.2f}). "
                "Some regions show inconsistent compression history."
            )
        else:
            score, severity = 0.12, SignalSeverity.LOW
            desc = (
                f"Natural ELA variation (mean={mean_ela:.1f}, CV={block_cv:.2f}). "
                "Consistent with real-world photography."
            )

        return score, DetectionSignal(
            name="Error Level Analysis (ELA)",
            description=desc,
            severity=severity,
            score=score,
            details=f"Mean={mean_ela:.2f} | Std={std_ela:.2f} | Block CV={block_cv:.3f} | Input={'JPEG' if is_jpeg else 'non-JPEG'}",
        )
    except Exception as exc:
        return 0.45, DetectionSignal(
            name="ELA Analysis Failed",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.45,
        )


def _noise_analysis(arr) -> Tuple[float, DetectionSignal]:
    """
    Camera sensors produce unique, near-white noise patterns.
    AI generators produce images with unnaturally smooth or correlated noise.
    """
    try:
        import numpy as np
        from scipy.ndimage import gaussian_filter

        gray   = arr.mean(axis=2)
        smooth = gaussian_filter(gray, sigma=2)
        noise  = gray - smooth

        noise_std = float(noise.std())

        h, w   = noise.shape
        sample = noise[:min(h, 256), :min(w, 256)].flatten()
        if len(sample) > 1:
            with np.errstate(invalid="ignore"):
                corr = np.corrcoef(sample[:-1], sample[1:])
            ac = float(corr[0, 1]) if not np.isnan(corr[0, 1]) else 0.0
        else:
            ac = 0.0

        if noise_std < 2.5:
            score, severity = 0.78, SignalSeverity.HIGH
            desc = f"Extremely low noise (σ={noise_std:.2f}). AI generators produce unrealistically clean images."
        elif noise_std < 5.0 and abs(ac) > 0.15:
            score, severity = 0.58, SignalSeverity.MEDIUM
            desc = f"Correlated noise (σ={noise_std:.2f}, AC={ac:.3f}). May indicate synthetic generation or heavy post-processing."
        elif abs(ac) > 0.25:
            score, severity = 0.48, SignalSeverity.MEDIUM
            desc = f"Spatially correlated noise (AC={ac:.3f}). Natural sensor noise is typically uncorrelated."
        else:
            score, severity = 0.12, SignalSeverity.LOW
            desc = f"Natural noise pattern (σ={noise_std:.2f}, AC={ac:.3f}). Consistent with real sensor noise."

        return score, DetectionSignal(
            name="Noise Pattern Analysis",
            description=desc,
            severity=severity,
            score=score,
            details=f"Noise σ={noise_std:.3f} | AC={ac:.4f}",
        )
    except Exception as exc:
        return 0.45, DetectionSignal(
            name="Noise Analysis Failed",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.45,
        )


def _frequency_analysis(arr) -> Tuple[float, DetectionSignal]:
    """
    DCT/FFT frequency domain analysis.
    Diffusion models produce images with anomalous high-frequency energy distribution.
    """
    try:
        import numpy as np

        gray = arr.mean(axis=2)
        h, w = gray.shape
        cy, cx = h // 2, w // 2
        size   = min(256, cy, cx)
        crop   = gray[cy - size:cy + size, cx - size:cx + size]

        fft_shift = np.fft.fftshift(np.fft.fft2(crop))
        magnitude = np.log1p(np.abs(fft_shift))

        mid_y, mid_x = magnitude.shape[0] // 2, magnitude.shape[1] // 2
        Y, X  = np.ogrid[:magnitude.shape[0], :magnitude.shape[1]]
        dist  = np.sqrt((Y - mid_y) ** 2 + (X - mid_x) ** 2)
        max_r = min(mid_y, mid_x)

        low_e  = float(magnitude[dist < max_r * 0.2].mean())
        mid_e  = float(magnitude[(dist >= max_r * 0.2) & (dist < max_r * 0.5)].mean())
        high_e = float(magnitude[dist >= max_r * 0.5].mean())

        ratio_hl = high_e / (low_e + 1e-6)

        if ratio_hl > 0.88:
            score, severity = 0.68, SignalSeverity.HIGH
            desc = f"Anomalous high-frequency energy (H/L={ratio_hl:.2f}). Diffusion models produce distinctive spectral signatures."
        elif ratio_hl > 0.72:
            score, severity = 0.42, SignalSeverity.MEDIUM
            desc = f"Elevated high-frequency energy (H/L={ratio_hl:.2f}). May indicate AI generation or heavy sharpening."
        else:
            score, severity = 0.12, SignalSeverity.LOW
            desc = f"Normal frequency distribution (H/L={ratio_hl:.2f}). Consistent with natural imagery."

        return score, DetectionSignal(
            name="DCT Frequency Analysis",
            description=desc,
            severity=severity,
            score=score,
            details=f"Low={low_e:.3f} | Mid={mid_e:.3f} | High={high_e:.3f} | H/L={ratio_hl:.3f}",
        )
    except Exception as exc:
        return 0.45, DetectionSignal(
            name="Frequency Analysis Failed",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.45,
        )


def _color_distribution(arr) -> Tuple[float, DetectionSignal]:
    """AI-generated images often exhibit over-saturated or hyper-consistent colour distributions."""
    try:
        import numpy as np

        r, g, b = arr[:, :, 0], arr[:, :, 1], arr[:, :, 2]
        r_std, g_std, b_std = float(r.std()), float(g.std()), float(b.std())

        rg_corr = float(np.corrcoef(r.flatten()[:10000], g.flatten()[:10000])[0, 1])
        mean_sat = float(_compute_saturation(arr).mean())

        if mean_sat > 0.78 and rg_corr > 0.95:
            score, severity = 0.62, SignalSeverity.MEDIUM
            desc = f"High saturation (S={mean_sat:.2f}) + channel correlation ({rg_corr:.2f}). AI images often over-saturate."
        elif r_std < 20 and g_std < 20 and b_std < 20:
            score, severity = 0.52, SignalSeverity.MEDIUM
            desc = f"Very low colour variance (R={r_std:.1f}, G={g_std:.1f}, B={b_std:.1f}). Unnaturally homogeneous."
        else:
            score, severity = 0.18, SignalSeverity.LOW
            desc = f"Normal colour distribution (saturation={mean_sat:.2f}). No anomalies detected."

        return score, DetectionSignal(
            name="Colour Distribution Analysis",
            description=desc,
            severity=severity,
            score=score,
            details=f"R σ={r_std:.1f} | G σ={g_std:.1f} | B σ={b_std:.1f} | Sat={mean_sat:.3f}",
        )
    except Exception as exc:
        return 0.45, DetectionSignal(
            name="Colour Analysis Failed",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.45,
        )


def _compute_saturation(arr):
    import numpy as np
    a = arr / 255.0
    c_max, c_min = a.max(axis=2), a.min(axis=2)
    return np.where(c_max > 0, (c_max - c_min) / (c_max + 1e-6), 0)


def _texture_regularity(arr) -> Tuple[float, DetectionSignal]:
    """AI images often have suspiciously uniform textures across the frame."""
    try:
        import numpy as np
        from scipy.ndimage import sobel

        gray = arr.mean(axis=2)
        h, w = gray.shape

        edges = np.sqrt(sobel(gray, axis=0) ** 2 + sobel(gray, axis=1) ** 2)
        block = 32
        densities = [
            float(edges[y:y + block, x:x + block].mean())
            for y in range(0, h - block, block)
            for x in range(0, w - block, block)
        ]
        d = np.array(densities)
        cv = float(d.std() / (d.mean() + 1e-6))

        if cv < 0.28:
            score, severity = 0.68, SignalSeverity.HIGH
            desc = f"Suspiciously regular texture (CV={cv:.3f}). AI images lack natural texture variation."
        elif cv < 0.52:
            score, severity = 0.38, SignalSeverity.MEDIUM
            desc = f"Moderately uniform texture (CV={cv:.3f})."
        else:
            score, severity = 0.12, SignalSeverity.LOW
            desc = f"Natural texture variation (CV={cv:.3f})."

        return score, DetectionSignal(
            name="Texture Regularity",
            description=desc,
            severity=severity,
            score=score,
            details=f"Edge density mean={d.mean():.2f} | std={d.std():.2f} | CV={cv:.4f}",
        )
    except Exception as exc:
        return 0.45, DetectionSignal(
            name="Texture Analysis Failed",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.45,
        )


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENT TEXT HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf_text(file_path: Path) -> str:
    import PyPDF2
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        return "\n".join(page.extract_text() or "" for page in reader.pages[:20])


def _pdf_package_signal(file_path: Path, text: str) -> Tuple[float, DetectionSignal]:
    """
    PDF-specific package analysis.
    Uses metadata, page structure, and export fingerprints that are invisible
    to plain text extraction.
    """
    score = 0.12
    findings: list[str] = []
    creator = ""
    producer = ""
    page_count = 0
    image_pages = 0
    fontless_pages = 0
    blank_text_pages = 0

    suspicious_markers = [
        "chatgpt", "openai", "claude", "anthropic", "gemini", "google ai",
        "copilot", "notebooklm", "perplexity", "grok",
        "reportlab", "weasyprint", "wkhtmltopdf", "python", "python-docx",
        "generated by", "auto-generated", "autogenerated",
    ]

    try:
        import PyPDF2

        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            info = reader.metadata or {}
            page_count = len(reader.pages)

            creator = str(info.get("/Creator", "") or "")
            producer = str(info.get("/Producer", "") or "")
            title = str(info.get("/Title", "") or "")
            author = str(info.get("/Author", "") or "")

            meta_blob = " ".join([creator, producer, title, author]).lower()
            for marker in suspicious_markers:
                if marker in meta_blob:
                    if marker in {"chatgpt", "openai", "claude", "anthropic", "gemini", "copilot"}:
                        score = max(score, 0.97)
                    elif marker in {"generated by", "auto-generated", "autogenerated"}:
                        score = max(score, 0.90)
                    else:
                        score = max(score, 0.72)
                    findings.append(f"metadata marker '{marker}'")

            if not any([creator.strip(), producer.strip(), title.strip(), author.strip()]):
                score = max(score, 0.52)
                findings.append("metadata fields largely empty")

            extracted_pages = 0
            for page in reader.pages[:20]:
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    extracted_pages += 1
                else:
                    blank_text_pages += 1

                resources = page.get("/Resources")
                if resources and hasattr(resources, "get_object"):
                    resources = resources.get_object()
                if not isinstance(resources, dict):
                    continue

                fonts = resources.get("/Font")
                if fonts and hasattr(fonts, "get_object"):
                    fonts = fonts.get_object()
                if not fonts:
                    fontless_pages += 1

                xobj = resources.get("/XObject")
                if xobj and hasattr(xobj, "get_object"):
                    xobj = xobj.get_object()
                if isinstance(xobj, dict):
                    has_image = False
                    for obj in xobj.values():
                        if hasattr(obj, "get_object"):
                            obj = obj.get_object()
                        if isinstance(obj, dict) and obj.get("/Subtype") == "/Image":
                            has_image = True
                            break
                    if has_image:
                        image_pages += 1

            if page_count >= 1 and image_pages == page_count and blank_text_pages == page_count:
                score = max(score, 0.76)
                findings.append("image-only PDF with no extractable text")
            elif page_count >= 2 and blank_text_pages >= max(2, page_count // 2):
                score = max(score, 0.62)
                findings.append("many pages lack extractable text")

            if page_count >= 1 and fontless_pages == page_count and blank_text_pages == page_count and image_pages == page_count:
                score = max(score, 0.78)
                findings.append("pages have no font resources despite visible content container")

            extracted_words = len(re.findall(r"\b\w+\b", text))
            if extracted_words >= 180 and page_count >= 2:
                words_per_page = extracted_words / page_count
                if 40 <= words_per_page <= 250:
                    score = max(score, 0.56)
                    findings.append("uniform medium-density multi-page layout")

        if not findings:
            return 0.12, DetectionSignal(
                name="PDF Package Metadata",
                description="PDF structure and metadata look broadly typical.",
                severity=SignalSeverity.LOW,
                score=0.12,
                details=f"Creator={creator or 'n/a'} | Producer={producer or 'n/a'} | Pages={page_count}",
            )

        if score >= 0.90:
            severity = SignalSeverity.HIGH
        elif score >= 0.65:
            severity = SignalSeverity.MEDIUM
        else:
            severity = SignalSeverity.LOW

        return score, DetectionSignal(
            name="PDF Package Metadata",
            description=(
                "Internal PDF metadata or page structure shows automation footprints: "
                + ", ".join(findings[:3]) + "."
            ),
            severity=severity,
            score=score,
            details=(
                f"Creator={creator or 'n/a'} | Producer={producer or 'n/a'} | Pages={page_count} | "
                f"ImagePages={image_pages} | FontlessPages={fontless_pages} | BlankTextPages={blank_text_pages}"
            ),
        )
    except Exception as exc:
        return 0.35, DetectionSignal(
            name="PDF Package Metadata",
            description="Could not inspect PDF package internals.",
            severity=SignalSeverity.LOW,
            score=0.35,
            details=str(exc),
        )


_DOCX_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}


def _read_docx_xml(file_path: Path, member: str) -> Optional[ET.Element]:
    try:
        with zipfile.ZipFile(file_path) as zf:
            return ET.fromstring(zf.read(member))
    except Exception:
        return None


def _extract_docx_text(file_path: Path) -> str:
    root = _read_docx_xml(file_path, "word/document.xml")
    if root is None:
        return ""

    paragraphs: list[str] = []
    for para in root.findall(".//w:p", _DOCX_NS):
        text = "".join(node.text or "" for node in para.findall(".//w:t", _DOCX_NS)).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_pptx_text(file_path: Path) -> tuple:
    """Returns (full_text, per_slide_texts) from a PPTX file."""
    from pptx import Presentation
    prs = Presentation(file_path)
    slide_texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slide_texts.append("\n".join(parts))
    return "\n\n".join(slide_texts), slide_texts


def _extract_xlsx_text(file_path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            row_vals = [
                str(c.value).strip()
                for c in row
                if c.value and isinstance(c.value, str) and str(c.value).strip()
            ]
            if row_vals:
                lines.append(" | ".join(row_vals))
    wb.close()
    return "\n".join(lines)


def _extract_csv_text(file_path: Path) -> str:
    import csv
    lines = []
    with open(file_path, newline="", encoding="utf-8-sig", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            parts = [cell.strip() for cell in row if cell.strip()]
            if parts:
                lines.append(" | ".join(parts))
    return "\n".join(lines)


def _docx_package_signal(file_path: Path, text: str) -> Tuple[float, DetectionSignal]:
    """
    DOCX-specific automation footprint analysis.
    Looks for package metadata and generation quirks that ordinary authored
    Word files usually do not exhibit.
    """
    signals = []
    score = 0.10
    creator = ""
    description = ""
    application = ""
    words = 0
    paragraphs = 0
    media_count = 0

    try:
        with zipfile.ZipFile(file_path) as zf:
            names = set(zf.namelist())
            media_count = sum(1 for name in names if name.startswith("word/media/"))

            core_root = ET.fromstring(zf.read("docProps/core.xml")) if "docProps/core.xml" in names else None
            app_root = ET.fromstring(zf.read("docProps/app.xml")) if "docProps/app.xml" in names else None

        if core_root is not None:
            values = {child.tag.split("}")[-1]: (child.text or "").strip() for child in core_root}
            creator = values.get("creator", "")
            description = values.get("description", "")
            created = values.get("created", "")
            modified = values.get("modified", "")
            revision = values.get("revision", "")

            generator_markers = " ".join([creator, description]).lower()
            if "python-docx" in generator_markers:
                score = max(score, 0.96)
                signals.append("python-docx generator metadata")
            elif any(marker in generator_markers for marker in ("generated by", "auto-generated", "autogenerated")):
                score = max(score, 0.84)
                signals.append("explicit generated-by metadata")

            if created and modified and created == modified and created.startswith("2013-12-23"):
                score = max(score, 0.82)
                signals.append("default python-docx template timestamp")

            if revision in ("", "0", "1") and len(text.split()) > 120:
                score = max(score, 0.66)
                signals.append("substantial content with first-revision metadata")

        if app_root is not None:
            values = {child.tag.split("}")[-1]: (child.text or "").strip() for child in app_root}
            application = values.get("Application", "")
            words = int(values.get("Words", "0") or 0)
            paragraphs = int(values.get("Paragraphs", "0") or 0)

            extracted_words = len(re.findall(r"\b\w+\b", text))
            if extracted_words >= 80 and words == 0:
                score = max(score, 0.76)
                signals.append("app.xml word count left at zero")
            if len(text.splitlines()) >= 12 and paragraphs == 0:
                score = max(score, 0.72)
                signals.append("app.xml paragraph count left at zero")

        if media_count >= 8 and len(text.split()) >= 180:
            score = max(score, 0.62)
            signals.append("heavily templated docx package with many embedded media assets")

        if not signals:
            return 0.12, DetectionSignal(
                name="DOCX Package Metadata",
                description="DOCX package metadata looks broadly typical.",
                severity=SignalSeverity.LOW,
                score=0.12,
                details=f"Creator={creator or 'n/a'} | Application={application or 'n/a'} | Media={media_count}",
            )

        if score >= 0.90:
            severity = SignalSeverity.HIGH
        elif score >= 0.65:
            severity = SignalSeverity.MEDIUM
        else:
            severity = SignalSeverity.LOW

        return score, DetectionSignal(
            name="DOCX Package Metadata",
            description=(
                "Internal Word package metadata shows automation footprints: "
                + ", ".join(signals[:3]) + "."
            ),
            severity=severity,
            score=score,
            details=(
                f"Creator={creator or 'n/a'} | Description={description or 'n/a'} | "
                f"Application={application or 'n/a'} | Words={words} | Paragraphs={paragraphs} | Media={media_count}"
            ),
        )
    except Exception as exc:
        return 0.35, DetectionSignal(
            name="DOCX Package Metadata",
            description="Could not inspect DOCX package internals.",
            severity=SignalSeverity.LOW,
            score=0.35,
            details=str(exc),
        )


def _split_into_units(text: str):
    """
    Split text into analysable units including bullet points and numbered lists,
    not just prose sentences ending in .!?
    """
    # Split on sentence endings AND newlines (for bullet-point structured docs)
    raw = re.split(r"[.!?]+|\n+", text)
    return [s.strip() for s in raw if 2 <= len(s.split()) <= 80]


def _text_burstiness(text: str) -> Tuple[float, DetectionSignal]:
    """
    Burstiness = variability in sentence / unit length.
    Human writing is "bursty"; AI text is uniformly structured.
    Also works on bullet-point CVs and structured documents.
    """
    units = _split_into_units(text)
    lengths = [len(s.split()) for s in units]

    if len(lengths) < 4:
        return 0.3, DetectionSignal(
            name="Sentence Burstiness",
            description="Too few units to analyse.",
            severity=SignalSeverity.LOW,
            score=0.3,
        )

    mean_l = sum(lengths) / len(lengths)
    std_l  = math.sqrt(sum((x - mean_l) ** 2 for x in lengths) / len(lengths))
    cv     = std_l / (mean_l + 1e-6)

    if cv < 0.25:
        score, severity = 0.75, SignalSeverity.HIGH
        desc = f"Low burstiness (CV={cv:.2f}). AI text has more uniform sentence lengths than human writing."
    elif cv < 0.40:
        score, severity = 0.48, SignalSeverity.MEDIUM
        desc = f"Moderate burstiness (CV={cv:.2f}). Somewhat uniform sentence structure."
    else:
        score, severity = 0.12, SignalSeverity.LOW
        desc = f"Natural burstiness (CV={cv:.2f}). Length variety consistent with human authorship."

    return score, DetectionSignal(
        name="Sentence Length Burstiness",
        description=desc,
        severity=severity,
        score=score,
        details=f"Mean={mean_l:.1f} words | Std={std_l:.1f} | CV={cv:.3f}",
    )


def _vocabulary_richness(text: str) -> Tuple[float, DetectionSignal]:
    """Type-Token Ratio — AI text reuses vocabulary more than humans."""
    words = re.findall(r"\b[a-z]+\b", text.lower())
    if len(words) < 50:
        return 0.3, DetectionSignal(
            name="Vocabulary Richness",
            description="Insufficient text for vocabulary analysis.",
            severity=SignalSeverity.LOW,
            score=0.3,
        )

    window = 100
    ttrs = [
        len(set(words[i:i + window])) / window
        for i in range(0, len(words) - window, window)
    ]
    mean_ttr = sum(ttrs) / len(ttrs) if ttrs else 0.5

    if mean_ttr > 0.75:
        score, severity = 0.12, SignalSeverity.LOW
        desc = f"High vocabulary richness (TTR={mean_ttr:.2f}). Diverse word choice → human authorship."
    elif mean_ttr > 0.60:
        score, severity = 0.32, SignalSeverity.LOW
        desc = f"Average vocabulary richness (TTR={mean_ttr:.2f})."
    else:
        score, severity = 0.65, SignalSeverity.MEDIUM
        desc = f"Low vocabulary richness (TTR={mean_ttr:.2f}). Repetitive word choice is a marker of AI text."

    return score, DetectionSignal(
        name="Vocabulary Richness (TTR)",
        description=desc,
        severity=severity,
        score=score,
        details=f"TTR={mean_ttr:.3f} over {len(ttrs)} windows",
    )


def _transition_phrases(text: str) -> Tuple[float, DetectionSignal]:
    """AI text over-uses specific transitional phrases and filler constructs."""
    AI_TRANSITIONS = [
        # ── Classic AI essay phrases ──
        "in conclusion", "furthermore", "moreover", "it is important to note",
        "it is worth noting", "in summary", "to summarize", "as a result",
        "on the other hand", "in addition", "it should be noted",
        "in this context", "needless to say", "it goes without saying",
        "first and foremost", "last but not least", "to begin with",
        "in light of", "this allows", "this enables", "this ensures",
        "thus ensuring", "plays a crucial role", "it is crucial",
        "it is essential", "plays a vital role", "key role", "delve into",
        "dive into", "at the end of the day", "moving forward",
        "having said that", "with that being said", "it's worth mentioning",
        "as previously mentioned", "as mentioned earlier",
        # ── CV / professional AI buzzwords ──
        "results-driven", "results driven", "detail-oriented", "detail oriented",
        "quick learner", "fast learner", "proven track record",
        "strong ability to", "growing interest in", "passionate about",
        "highly motivated", "self-motivated", "self motivated",
        "team player", "go-getter", "fast-paced environment",
        "seeking to leverage", "seeking to utilize", "adept at",
        "proficient in", "well-versed in", "versed in",
        "dedicated to", "committed to", "striving to",
        "ensuring confidentiality", "attention to detail",
        "strong communication skills", "excellent communication",
        "analytical thinking", "analytical skills",
        "problem-solving skills", "problem solving skills",
        "data-driven", "data driven", "cross-functional",
        "stakeholder", "synergize", "leverage", "value-added",
        "actionable insights", "thought leader", "best practices",
        "continuous improvement", "proactive approach",
        "comprehensive understanding", "demonstrated ability",
        "with experience in", "experience in",
        "dynamic professional", "dynamic individual",
        "innovative solutions", "strategic thinking",
        "driving results", "impactful", "collaborative",
        # ── Structural AI giveaways ──
        "comprehensive", "multifaceted", "robust", "streamline",
        "cutting-edge", "state-of-the-art", "holistic approach",
        "tailor", "tailored", "bespoke", "seamlessly",
        "transformative", "empower", "foster", "facilitate",
        "spearhead", "spearheaded", "championed",
    ]
    lower      = text.lower()
    word_count = max(len(text.split()), 1)
    found      = [(p, lower.count(p)) for p in AI_TRANSITIONS if p in lower]
    total_hits = sum(c for _, c in found)
    density    = total_hits / (word_count / 100)

    if density > 2.0:
        score, severity = 0.80, SignalSeverity.HIGH
        top  = sorted(found, key=lambda x: x[1], reverse=True)[:5]
        desc = f"High AI phrase density ({density:.1f}/100 words). Top phrases: {', '.join(p for p, _ in top)}"
    elif density > 1.0:
        score, severity = 0.48, SignalSeverity.MEDIUM
        desc = f"Moderate AI phrase density ({density:.1f}/100 words)."
    else:
        score, severity = 0.08, SignalSeverity.LOW
        desc = f"Low AI phrase density ({density:.1f}/100 words)."

    return score, DetectionSignal(
        name="AI Transition Phrase Density",
        description=desc,
        severity=severity,
        score=score,
        details=f"Density={density:.2f}/100 words | Unique phrases: {len(found)}",
    )


def _paragraph_homogeneity(text: str) -> Tuple[float, DetectionSignal]:
    """AI tends to write paragraphs of very similar length."""
    # Try double-newline paragraphs first; fall back to single-newline blocks
    # so structured docs (CVs, reports) that use \n instead of \n\n still work.
    paragraphs = [p.strip() for p in text.split("\n\n") if len(p.strip().split()) > 10]
    if len(paragraphs) < 3:
        # Single-newline fallback for CVs and structured docs (bullet points are shorter)
        paragraphs = [p.strip() for p in text.split("\n") if len(p.strip().split()) > 5]
    if len(paragraphs) < 3:
        return 0.3, DetectionSignal(
            name="Paragraph Homogeneity",
            description="Too few paragraphs to assess.",
            severity=SignalSeverity.LOW,
            score=0.3,
        )

    lengths = [len(p.split()) for p in paragraphs]
    mean_l  = sum(lengths) / len(lengths)
    std_l   = math.sqrt(sum((x - mean_l) ** 2 for x in lengths) / len(lengths))
    cv      = std_l / (mean_l + 1e-6)

    if cv < 0.20:
        score, severity = 0.70, SignalSeverity.HIGH
        desc = f"Highly uniform paragraph lengths (CV={cv:.2f}). AI produces paragraphs of near-identical length."
    elif cv < 0.40:
        score, severity = 0.38, SignalSeverity.MEDIUM
        desc = f"Somewhat uniform paragraphs (CV={cv:.2f})."
    else:
        score, severity = 0.12, SignalSeverity.LOW
        desc = f"Natural paragraph length variation (CV={cv:.2f})."

    return score, DetectionSignal(
        name="Paragraph Length Homogeneity",
        description=desc,
        severity=severity,
        score=score,
        details=f"Mean={mean_l:.0f} words | Std={std_l:.1f} | CV={cv:.3f}",
    )


def _sentence_self_similarity(text: str) -> Tuple[float, DetectionSignal]:
    """
    Measures lexical overlap between consecutive sentences using Jaccard similarity.
    AI text tends to stay on the same narrow vocabulary; human text shifts more.
    """
    sentences = [s for s in _split_into_units(text) if len(s.split()) > 3]
    if len(sentences) < 5:
        return 0.3, DetectionSignal(
            name="Sentence Self-Similarity",
            description="Too few sentences for self-similarity analysis.",
            severity=SignalSeverity.LOW,
            score=0.3,
        )

    vectors = [set(s.lower().split()) for s in sentences]
    sims = []
    for i in range(len(vectors)):
        for j in range(i + 1, min(i + 4, len(vectors))):
            a, b = vectors[i], vectors[j]
            union = a | b
            if union:
                sims.append(len(a & b) / len(union))

    mean_sim = sum(sims) / len(sims) if sims else 0.0

    if mean_sim > 0.22:
        score, severity = 0.65, SignalSeverity.MEDIUM
        desc = f"High sentence self-similarity ({mean_sim:.2f}). AI text reuses the same vocabulary within paragraphs."
    elif mean_sim > 0.12:
        score, severity = 0.35, SignalSeverity.LOW
        desc = f"Moderate sentence self-similarity ({mean_sim:.2f})."
    else:
        score, severity = 0.10, SignalSeverity.LOW
        desc = f"Low sentence self-similarity ({mean_sim:.2f}). Vocabulary shifts naturally — consistent with human writing."

    return score, DetectionSignal(
        name="Sentence Self-Similarity",
        description=desc,
        severity=severity,
        score=score,
        details=f"Mean Jaccard similarity={mean_sim:.3f} over {len(sims)} pairs",
    )


def _perplexity_signal(
    text: str,
    signals: List[DetectionSignal],
) -> Tuple[Optional[float], DetectionSignal]:
    """GPT-2 perplexity — lower = more predictable = more AI-like."""
    try:
        from app.services.ml_detector import ml_perplexity_score
        prob, details = ml_perplexity_score(text)
        if prob is None:
            sig = DetectionSignal(
                name="GPT-2 Perplexity",
                description="Perplexity analysis unavailable.",
                severity=SignalSeverity.LOW,
                score=0.45,
                details=details,
            )
            return None, sig

        if prob >= 0.75:
            severity = SignalSeverity.HIGH
        elif prob >= 0.50:
            severity = SignalSeverity.MEDIUM
        else:
            severity = SignalSeverity.LOW

        sig = DetectionSignal(
            name="GPT-2 Perplexity",
            description=(
                f"AI-generated text is statistically more predictable under language models. "
                f"This text scores {prob:.0%} AI-likelihood based on predictability."
            ),
            severity=severity,
            score=prob,
            details=details,
        )
        return prob, sig
    except Exception as exc:
        logger.warning("Perplexity signal failed: %s", exc)
        return None, DetectionSignal(
            name="GPT-2 Perplexity",
            description="Perplexity analysis failed.",
            severity=SignalSeverity.LOW,
            score=0.45,
        )


def _template_placeholder_check(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect template placeholders and generic fill-in-the-blank patterns.
    Brackets like [Your Name], 'Company Name', 'University Name' are
    strong indicators of AI-generated or template content.
    """
    PLACEHOLDER_PATTERNS = [
        r"\[your [^\]]{1,40}\]",        # [Your Full Name], [Your Email] …
        r"\[name\]", r"\[email\]",
        r"\[company[^\]]{0,20}\]",
        r"\[job title\]", r"\[position\]",
        r"\[university[^\]]{0,20}\]",
        r"\[degree[^\]]{0,20}\]",
        r"\[year[^\]]{0,10}\]",
        r"\byour full name\b",
        r"\bcompany name\b",
        r"\buniversity name\b",
        r"\bjob title\b\s*[-–]",
        r"\bdegree\b\s*[-–]\s*university",
        r"\bcity,\s*country\b",
        r"\(year\s*[-–]\s*present\)",
        r"\(year\)",
        r"\byour name\b",
        r"\byour email\b",
        r"\byour phone\b",
    ]
    lower = text.lower()
    found = []
    for pat in PLACEHOLDER_PATTERNS:
        matches = re.findall(pat, lower)
        found.extend(matches)

    count = len(found)
    if count >= 3:
        score, severity = 0.96, SignalSeverity.HIGH
        desc = (
            f"Found {count} template placeholders (e.g. {', '.join(repr(f) for f in found[:3])}). "
            "These are strong markers of AI-generated or template content."
        )
    elif count >= 1:
        score, severity = 0.82, SignalSeverity.HIGH
        desc = f"Found template placeholder: {', '.join(repr(f) for f in found)}. Indicates AI-generated or template content."
    else:
        score, severity = 0.05, SignalSeverity.LOW
        desc = "No template placeholders detected."

    return score, DetectionSignal(
        name="Template Placeholders",
        description=desc,
        severity=severity,
        score=score,
        details=f"Placeholders found: {count}",
    )


def _professional_buzzwords(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect high-density professional buzzwords that AI models overuse
    in CVs, cover letters, and business documents.
    """
    # Note: terms already in AI_TRANSITIONS are intentionally omitted here
    # to avoid double-counting the same signal across both checks.
    BUZZWORDS = [
        "results-driven", "results driven", "detail-oriented", "detail oriented",
        "quick learner", "fast learner", "proven track record",
        "strong ability", "growing interest", "passionate about",
        "highly motivated", "self-motivated", "team player",
        "fast-paced", "seeking to leverage", "adept at",
        "dedicated professional", "committed to excellence",
        "ensuring confidentiality", "data accuracy",
        "strong communication", "excellent communication",
        "analytical thinking", "problem-solving",
        "actionable insights",
        "demonstrated ability",
        "dynamic professional", "innovative solutions",
        "strategic thinking", "driving results",
        "holistic approach",
        "leveraging", "value-added", "collaborative environment",
        "attention to detail", "time management",
        "communication skills", "teamwork",
        "go-getter", "outside the box", "think outside",
        "hard-working", "hardworking", "reliable professional",
        "exceed expectations", "above and beyond",
    ]
    lower = text.lower()
    word_count = max(len(text.split()), 1)
    found = [(bw, lower.count(bw)) for bw in BUZZWORDS if bw in lower]
    total_hits = sum(c for _, c in found)
    density = total_hits / (word_count / 100)

    if density > 4.0 or len(found) >= 6:
        score, severity = 0.88, SignalSeverity.HIGH
        top = sorted(found, key=lambda x: x[1], reverse=True)[:5]
        desc = (
            f"Very high professional buzzword density ({density:.1f}/100 words, {len(found)} unique). "
            f"Top: {', '.join(bw for bw, _ in top)}. "
            "AI models saturate professional documents with these phrases."
        )
    elif density > 2.0 or len(found) >= 3:
        score, severity = 0.65, SignalSeverity.MEDIUM
        desc = f"Elevated buzzword density ({density:.1f}/100 words, {len(found)} unique). Common in AI-written professional content."
    elif len(found) >= 1:
        score, severity = 0.35, SignalSeverity.LOW
        desc = f"Low buzzword density ({density:.1f}/100 words). A few common professional phrases detected."
    else:
        score, severity = 0.08, SignalSeverity.LOW
        desc = "No notable professional AI buzzwords detected."

    return score, DetectionSignal(
        name="Professional AI Buzzwords",
        description=desc,
        severity=severity,
        score=score,
        details=f"Density={density:.2f}/100 words | Unique buzzwords: {len(found)}",
    )


# ─────────────────────────────────────────────────────────────────────────────
# HR-SPECIFIC DETECTION SIGNALS
# ─────────────────────────────────────────────────────────────────────────────

def _cv_section_structure(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect rigid AI CV template section headers.
    AI-generated CVs follow a predictable, identical section pattern every time.
    """
    STANDARD_SECTIONS = [
        r"\bprofessional summary\b", r"\bcareer objective\b", r"\bobjective\b",
        r"\bwork experience\b", r"\bprofessional experience\b", r"\bemployment history\b",
        r"\beducation\b", r"\bskills\b", r"\bcore competencies\b",
        r"\bcertifications?\b", r"\bcertificates?\b", r"\bvolunteer\b",
        r"\bextra.?curricular\b", r"\bhobbies\b", r"\breferences\b",
        r"\bpublications?\b", r"\bawards?\b", r"\bkey skills\b",
        r"\btechnical skills\b", r"\blanguages\b", r"\bachievements?\b",
        r"\bprofile\b", r"\bsummary\b", r"\binterests?\b",
    ]
    lower = text.lower()
    found = [p for p in STANDARD_SECTIONS if re.search(p, lower)]
    count = len(found)

    if count >= 6:
        score, severity = 0.85, SignalSeverity.HIGH
        desc = (
            f"Detected {count} standard CV section headers. AI tools produce rigidly "
            "templated CVs with identical section structures — Professional Summary, "
            "Work Experience, Education, Skills, Certifications, References."
        )
    elif count >= 4:
        score, severity = 0.65, SignalSeverity.MEDIUM
        desc = f"Detected {count} standard CV/document section headers. Moderately templated structure."
    elif count >= 2:
        score, severity = 0.35, SignalSeverity.LOW
        desc = f"Detected {count} standard section headers. Low structural signal."
    else:
        score, severity = 0.05, SignalSeverity.LOW
        desc = "No rigid CV/document section structure detected."

    return score, DetectionSignal(
        name="CV Section Structure",
        description=desc,
        severity=severity,
        score=score,
        details=f"Standard section headers found: {count}",
    )


def _achievement_claim_density(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect vague quantified achievement claims.
    AI models fill CVs with impressive-sounding but unverifiable percentage metrics.
    """
    ACHIEVEMENT_PATTERNS = [
        r"\b(?:increased|improved|boosted|grew|enhanced|raised)\b.{0,50}\b\d+\s*%",
        r"\b(?:reduced|decreased|cut|lowered|minimized|minimised)\b.{0,50}\b\d+\s*%",
        r"\b(?:saved|generated|drove|delivered|achieved)\b.{0,50}\b\d+\s*%",
        r"\b(?:exceeded|surpassed|outperformed)\b.{0,50}\b(?:target|goal|kpi|quota)",
        r"\b\d+\s*%\s+(?:increase|improvement|reduction|growth|efficiency)",
        r"\btop\s+\d+\s*%\b",
        r"\bover\s+\d+\s+(?:years?|months?)\s+of\s+experience\b",
        r"\b(?:successfully|consistently)\s+(?:delivered|met|exceeded|achieved)\b",
        r"\bkey\s+(?:contributor|player|driver)\b",
        r"\bresulted?\s+in\s+(?:a\s+)?\d+",
        r"\b(?:spearheaded|championed|orchestrated)\b.{0,40}\b(?:leading|resulting|achieving)",
        r"\bimpact(?:ed|ful)\b.{0,30}\b(?:\d+|\bmillion\b|\bthousand\b)",
    ]
    lower = text.lower()
    hits = []
    for pat in ACHIEVEMENT_PATTERNS:
        hits.extend(re.findall(pat, lower))

    count = len(hits)
    word_count = max(len(text.split()), 1)
    density = count / (word_count / 100)

    if count >= 5 or density > 2.0:
        score, severity = 0.82, SignalSeverity.HIGH
        desc = (
            f"High density of vague achievement claims ({count} found, {density:.1f}/100 words). "
            "AI models saturate CVs with unverifiable metrics: 'increased efficiency by 30%', "
            "'reduced costs by 25%', 'exceeded KPIs consistently'."
        )
    elif count >= 2 or density > 0.8:
        score, severity = 0.55, SignalSeverity.MEDIUM
        desc = f"Moderate achievement claim density ({count} claims). Common pattern in AI-generated professional documents."
    elif count >= 1:
        score, severity = 0.25, SignalSeverity.LOW
        desc = f"Low achievement claim density ({count} claims found)."
    else:
        score, severity = 0.05, SignalSeverity.LOW
        desc = "No suspicious vague achievement claim patterns detected."

    return score, DetectionSignal(
        name="Achievement Claim Density",
        description=desc,
        severity=severity,
        score=score,
        details=f"Claims found: {count} | Density: {density:.2f}/100 words",
    )


def _skill_list_density(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect unrealistically long skill/technology lists.
    AI-generated CVs enumerate an implausibly large and diverse set of skills.
    """
    TECH_SEPARATORS = re.compile(r"[,|•·▪\-–—/]")
    SKILL_PATTERN = re.compile(
        r"\b(?:python|java(?:script)?|typescript|react|angular|vue|node|sql|"
        r"excel|powerpoint|word|sap|salesforce|jira|git|aws|azure|gcp|"
        r"tableau|power\s*bi|agile|scrum|kanban|lean|six\s*sigma|"
        r"c\+\+|c#|\.net|php|ruby|swift|kotlin|docker|kubernetes|"
        r"machine\s*learning|deep\s*learning|nlp|data\s*analysis|"
        r"communication|leadership|management|microsoft\s*office|"
        r"google\s*workspace|crm|erp|html|css|rest\s*api|ci/cd|devops)\b",
        re.IGNORECASE,
    )

    lines = text.split("\n")
    total_skill_items = 0
    max_skills_line = 0

    for line in lines:
        items = [s.strip() for s in TECH_SEPARATORS.split(line) if s.strip()]
        if len(items) >= 3:
            skill_count = sum(
                1 for item in items
                if SKILL_PATTERN.search(item) or len(item.split()) <= 3
            )
            if skill_count >= 3:
                total_skill_items += skill_count
                max_skills_line = max(max_skills_line, skill_count)

    skill_keyword_count = len(SKILL_PATTERN.findall(text))

    if total_skill_items >= 15 or skill_keyword_count >= 20:
        score, severity = 0.80, SignalSeverity.HIGH
        desc = (
            f"Implausibly dense skill enumeration ({total_skill_items} listed items, "
            f"{skill_keyword_count} technology keywords). AI-generated CVs list an "
            "unrealistically broad range of skills to appear more qualified."
        )
    elif total_skill_items >= 8 or skill_keyword_count >= 12:
        score, severity = 0.55, SignalSeverity.MEDIUM
        desc = f"High skill density ({total_skill_items} list items, {skill_keyword_count} keywords). Common in AI-written CVs."
    elif total_skill_items >= 4 or skill_keyword_count >= 6:
        score, severity = 0.28, SignalSeverity.LOW
        desc = f"Moderate skill density ({total_skill_items} items, {skill_keyword_count} keywords)."
    else:
        score, severity = 0.05, SignalSeverity.LOW
        desc = "Normal skill enumeration. No unusual density detected."

    return score, DetectionSignal(
        name="Skill List Density",
        description=desc,
        severity=severity,
        score=score,
        details=f"Listed items: {total_skill_items} | Skill keywords: {skill_keyword_count} | Max per line: {max_skills_line}",
    )


def _slide_structural_uniformity(slide_texts: list) -> Tuple[float, DetectionSignal]:
    """
    Detect AI presentations by per-slide word count uniformity.
    AI-generated presentations produce slides with nearly identical content lengths.
    """
    lengths = [len(s.split()) for s in slide_texts if s.strip()]
    if len(lengths) < 3:
        return 0.30, DetectionSignal(
            name="Slide Structure Uniformity",
            description="Too few slides to assess uniformity.",
            severity=SignalSeverity.LOW,
            score=0.30,
        )

    mean_l = sum(lengths) / len(lengths)
    std_l = math.sqrt(sum((x - mean_l) ** 2 for x in lengths) / len(lengths))
    cv = std_l / (mean_l + 1e-6)

    AI_SLIDE_TITLES = [
        "introduction", "overview", "background", "agenda", "conclusion",
        "summary", "thank you", "questions", "key takeaways", "next steps",
        "recommendations", "objectives", "scope",
    ]
    first_words = []
    for s in slide_texts[:4]:
        words = s.strip().split()
        if words:
            first_words.append(words[0].lower().rstrip(":"))
    title_hits = sum(1 for w in first_words if any(t.startswith(w) or w.startswith(t) for t in AI_SLIDE_TITLES))

    if cv < 0.25:
        score, severity = 0.80, SignalSeverity.HIGH
        desc = (
            f"Highly uniform slide content lengths (CV={cv:.2f}, mean={mean_l:.0f} words/slide). "
            "AI-generated presentations produce slides with near-identical word counts."
        )
    elif cv < 0.45:
        score, severity = 0.48, SignalSeverity.MEDIUM
        desc = f"Moderately uniform slide structure (CV={cv:.2f})."
    else:
        score, severity = 0.15, SignalSeverity.LOW
        desc = f"Natural slide length variation (CV={cv:.2f})."

    if title_hits >= 2:
        score = min(score + 0.12, 0.92)
        desc += f" Generic AI slide titles detected in opening slides ({title_hits} matches)."
        if severity == SignalSeverity.LOW:
            severity = SignalSeverity.MEDIUM

    return score, DetectionSignal(
        name="Slide Structure Uniformity",
        description=desc,
        severity=severity,
        score=score,
        details=f"Slides: {len(lengths)} | Mean words/slide: {mean_l:.0f} | CV: {cv:.3f} | Title hits: {title_hits}",
    )


def _cover_letter_patterns(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect AI-generated professional prose patterns in cover letters and summaries.
    These specific constructions are heavily over-represented in AI-generated professional text.
    """
    PATTERNS = [
        r"\bI am (?:writing|excited|thrilled|delighted|pleased) to (?:apply|express|submit|introduce)",
        r"\bI am a (?:highly|results|detail|passionate|dedicated|motivated|self-motivated|driven)",
        r"\bhave (?:a )?(?:strong|proven|solid|extensive|deep|broad) (?:background|experience|expertise|foundation) in",
        r"\bI would (?:welcome|value|appreciate|love) the opportunity",
        r"\bI look forward to (?:discussing|hearing from|speaking with|the opportunity|connecting)",
        r"\bthank you for (?:your|the) (?:time|consideration|opportunity|review)",
        r"\bI am confident (?:that|in my|I can|I will)",
        r"\bmy (?:strong|extensive|solid|diverse|broad|unique) (?:background|experience|skill set|expertise|qualifications)",
        r"\bpassion(?:ate)? (?:for|about) (?:this|the) (?:role|position|field|industry|opportunity|company)",
        r"\bI (?:possess|bring|offer) (?:a|an|strong|extensive|unique|diverse)",
        r"\bI am eager to (?:contribute|join|bring|apply|learn|grow|leverage)",
        r"\bmy (?:commitment|dedication|drive|passion|enthusiasm) to",
        r"\bseek(?:ing)? to (?:contribute|leverage|utilize|apply|join|grow)",
        r"\bopportunity to (?:contribute|join|grow|develop|work with|be part)",
        r"\bwith (?:my|a|an|over) (?:\d+ )?(?:years?|strong|extensive|proven) (?:of )?(?:experience|background)",
        r"\bI (?:successfully|consistently|effectively|efficiently) (?:managed|led|delivered|achieved|developed|implemented)",
        r"\bwould be (?:a )?(?:great|perfect|excellent|ideal|strong|valuable) (?:fit|match|addition|asset)",
        r"\bI am (?:particularly|especially|genuinely|truly|deeply) (?:excited|interested|passionate|motivated) (?:about|by|in)",
        r"\bI have (?:always|long) been (?:passionate|fascinated|interested|motivated)",
        r"\bI am a (?:highly|results|goal|customer|team|data|detail|self)-\w+",
    ]
    lower = text.lower()
    found_count = sum(1 for pat in PATTERNS if re.search(pat, lower))
    word_count = max(len(text.split()), 1)
    density = found_count / (word_count / 100)

    if found_count >= 5 or density > 1.5:
        score, severity = 0.88, SignalSeverity.HIGH
        desc = (
            f"Heavy use of AI cover letter templates ({found_count} patterns). "
            "Phrases like 'I am passionate about', 'I look forward to', 'I am confident that' "
            "are characteristic of AI-generated professional prose."
        )
    elif found_count >= 3 or density > 0.8:
        score, severity = 0.72, SignalSeverity.HIGH
        desc = f"Multiple AI professional prose patterns detected ({found_count}). Common in AI-generated cover letters and summaries."
    elif found_count >= 1 or density > 0.3:
        score, severity = 0.42, SignalSeverity.MEDIUM
        desc = f"Some AI professional prose patterns detected ({found_count})."
    else:
        score, severity = 0.05, SignalSeverity.LOW
        desc = "No characteristic AI cover letter patterns detected."

    return score, DetectionSignal(
        name="AI Professional Prose Patterns",
        description=desc,
        severity=severity,
        score=score,
        details=f"Patterns matched: {found_count} | Density: {density:.2f}/100 words",
    )


def _action_verb_density(text: str) -> Tuple[float, DetectionSignal]:
    """
    Detect AI CV bullet point patterns.
    AI-generated CVs consistently open bullet points with the same narrow set of action verbs.
    """
    AI_ACTION_VERBS = {
        "managed", "led", "developed", "implemented", "coordinated", "collaborated",
        "achieved", "spearheaded", "ensured", "facilitated", "oversaw", "executed",
        "delivered", "streamlined", "optimized", "leveraged", "drove", "championed",
        "orchestrated", "established", "maintained", "supported", "utilized",
        "conducted", "created", "designed", "built", "improved", "enhanced",
        "increased", "reduced", "analyzed", "evaluated", "identified", "provided",
        "prepared", "presented", "communicated", "fostered", "generated",
        "contributed", "overseeing", "leading", "managing", "developing",
        "implementing", "liaised", "liaising", "monitored", "ensured",
    }
    lines = text.split("\n")
    bullet_starts = []
    total_content_lines = 0
    for line in lines:
        stripped = line.strip().lstrip("•·▪-–—*>0123456789.").strip()
        if stripped and len(stripped.split()) >= 3:
            total_content_lines += 1
            first_word = stripped.split()[0].lower().rstrip(".,;:")
            if first_word in AI_ACTION_VERBS:
                bullet_starts.append(first_word)

    if total_content_lines < 4:
        return 0.25, DetectionSignal(
            name="Action Verb Patterns",
            description="Too few lines for action verb analysis.",
            severity=SignalSeverity.LOW,
            score=0.25,
        )

    ratio = len(bullet_starts) / max(total_content_lines, 1)
    unique_verbs = len(set(bullet_starts))

    if ratio >= 0.45 and unique_verbs >= 5:
        score, severity = 0.82, SignalSeverity.HIGH
        top = sorted(set(bullet_starts), key=lambda v: bullet_starts.count(v), reverse=True)[:5]
        desc = (
            f"Heavy use of AI-typical action verbs ({ratio:.0%} of lines, {unique_verbs} unique). "
            f"Top: {', '.join(top)}. AI CVs consistently open bullets with the same action verbs."
        )
    elif ratio >= 0.30 and unique_verbs >= 3:
        score, severity = 0.55, SignalSeverity.MEDIUM
        desc = f"Elevated action verb dependency ({ratio:.0%} of lines use AI-typical verb starters)."
    elif ratio >= 0.15:
        score, severity = 0.25, SignalSeverity.LOW
        desc = f"Moderate action verb usage ({ratio:.0%} of lines)."
    else:
        score, severity = 0.08, SignalSeverity.LOW
        desc = "Low action verb repetition. Natural writing variation detected."

    return score, DetectionSignal(
        name="Action Verb Patterns",
        description=desc,
        severity=severity,
        score=score,
        details=f"Ratio: {ratio:.1%} | Unique verbs: {unique_verbs} | Action starts: {len(bullet_starts)}/{total_content_lines}",
    )


_COMMON_ENGLISH = {
    "the", "and", "to", "of", "a", "in", "is", "it", "you", "that", "he",
    "was", "for", "on", "are", "as", "with", "his", "they", "at", "be",
    "this", "from", "or", "an", "by", "have", "not", "but", "we", "what",
    "which", "one", "had", "your", "can", "has", "will", "more", "if",
    "about", "up", "out", "do", "their", "all", "would", "there", "when",
    "use", "any", "how", "each", "she", "her", "than", "then", "its", "my",
    "into", "been", "also", "through", "time", "no", "so", "our", "new",
    "who", "me", "should", "just", "over", "these", "may", "after", "first",
    "work", "well", "way", "even", "only", "such", "much", "because", "too",
}


def _english_word_ratio(text: str) -> float:
    """Returns fraction of tokens that are common English function words."""
    words = re.findall(r"\b[a-z]+\b", text.lower())
    if len(words) < 20:
        return 1.0
    return sum(1 for w in words if w in _COMMON_ENGLISH) / len(words)


def _document_section_uniformity(text: str) -> Tuple[float, DetectionSignal]:
    """
    Language-agnostic structural uniformity detector.
    AI generators produce documents with near-identical section sizes.
    Works on workout plans, lesson plans, guides — anything structured.
    """
    blocks = [b.strip() for b in re.split(r"\n{2,}", text) if b.strip() and len(b.split()) >= 8]
    if len(blocks) < 3:
        blocks = [b.strip() for b in text.split("\n") if b.strip() and len(b.split()) >= 8]

    if len(blocks) < 3:
        return 0.25, DetectionSignal(
            name="Section Structure Uniformity",
            description="Too few sections to assess structural uniformity.",
            severity=SignalSeverity.LOW,
            score=0.25,
        )

    lengths = [len(b.split()) for b in blocks]
    mean_l = sum(lengths) / len(lengths)
    std_l = math.sqrt(sum((x - mean_l) ** 2 for x in lengths) / len(lengths))
    cv = std_l / (mean_l + 1e-6)

    if cv < 0.18:
        score, severity = 0.82, SignalSeverity.HIGH
        desc = (
            f"Highly uniform section lengths (CV={cv:.2f}, mean={mean_l:.0f} words/block, "
            f"{len(blocks)} sections). AI-generated structured documents have near-identical section sizes."
        )
    elif cv < 0.32:
        score, severity = 0.55, SignalSeverity.MEDIUM
        desc = f"Moderately uniform section lengths (CV={cv:.2f}). Some structural repetition detected."
    elif cv < 0.50:
        score, severity = 0.28, SignalSeverity.LOW
        desc = f"Slight uniformity in section lengths (CV={cv:.2f})."
    else:
        score, severity = 0.10, SignalSeverity.LOW
        desc = f"Natural section length variation (CV={cv:.2f})."

    return score, DetectionSignal(
        name="Section Structure Uniformity",
        description=desc,
        severity=severity,
        score=score,
        details=f"Blocks: {len(blocks)} | Mean words: {mean_l:.0f} | CV: {cv:.3f}",
    )


# ─────────────────────────────────────────────────────────────────────────────
# CORE TEXT ANALYSIS PIPELINE (shared by document, presentation, spreadsheet)
# ─────────────────────────────────────────────────────────────────────────────

def _analyse_text_content(text: str) -> Tuple[float, List[DetectionSignal]]:
    """
    Full multi-signal AI text detection pipeline.
    Used by analyze_document, analyze_presentation, and analyze_spreadsheet.
    """
    signals: List[DetectionSignal] = []

    if not text.strip():
        signals.append(DetectionSignal(
            name="Empty Content",
            description="No extractable text found in this file.",
            severity=SignalSeverity.LOW,
            score=0.0,
        ))
        return 0.3, signals

    # Detect document language early — determines which signal path to use
    eng_ratio = _english_word_ratio(text)
    is_english = eng_ratio > 0.07  # < 7% common English function words = non-English

    # ── Language-agnostic structural signal (always computed) ─────────────
    section_score, section_sig = _document_section_uniformity(text)

    # ── Language-agnostic signals (reliable for any language) ─────────────
    burst_score,  burst_sig  = _text_burstiness(text)
    vocab_score,  vocab_sig  = _vocabulary_richness(text)
    para_score,   para_sig   = _paragraph_homogeneity(text)
    sim_score,    sim_sig    = _sentence_self_similarity(text)

    if not is_english:
        # Non-English document: English ML models and pattern signals are unreliable.
        # Use only language-agnostic structural signals.
        signals += [section_sig, burst_sig, vocab_sig, para_sig, sim_sig]
        signals.append(DetectionSignal(
            name="Non-English Document Detected",
            description=(
                f"Document is not in English (English word ratio: {eng_ratio:.1%}). "
                "English-specific pattern signals cannot be applied. "
                "Structural uniformity and statistical signals used. "
                "Manual review recommended for non-English documents."
            ),
            severity=SignalSeverity.MEDIUM,
            score=0.5,
            details=f"English word ratio: {eng_ratio:.3f}",
        ))

        # Language-agnostic weighted scoring
        lang_scores   = [section_score, burst_score, vocab_score, para_score, sim_score]
        lang_weights  = [0.35,          0.28,        0.15,        0.14,       0.08]  # sum=1.00
        h_prob = sum(s * w for s, w in zip(lang_scores, lang_weights))

        # Base floor: cannot confidently clear non-English documents — default to "Uncertain"
        ai_prob = max(h_prob, 0.42)

        # Consensus floor for language-agnostic signals
        high_lang = sum(1 for s in lang_scores if s >= 0.65)
        if high_lang >= 3:
            ai_prob = max(ai_prob, 0.82)
        elif high_lang >= 2:
            ai_prob = max(ai_prob, 0.70)
        elif high_lang >= 1:
            ai_prob = max(ai_prob, 0.58)

        # Section uniformity + burstiness combo flags structured AI content
        if section_score >= 0.55 and burst_score >= 0.45:
            ai_prob = max(ai_prob, 0.65)
        if section_score >= 0.75 and burst_score >= 0.55:
            ai_prob = max(ai_prob, 0.75)

        return float(min(max(ai_prob, 0.0), 1.0)), signals

    # ── English path: full signal suite ───────────────────────────────────

    # ── 1. RoBERTa ML classifier ───────────────────────────────────────────
    ml_prob, _ml_weight = _ml_text_signal(text, signals)

    # ── 2. GPT-2 perplexity ────────────────────────────────────────────────
    perp_score, perp_sig = _perplexity_signal(text, signals)

    # ── 3. Template placeholders ──────────────────────────────────────────
    tmpl_score, tmpl_sig = _template_placeholder_check(text)

    # ── 4. Professional AI buzzwords ──────────────────────────────────────
    buzz_score, buzz_sig = _professional_buzzwords(text)

    # ── 5. CV / document section structure ────────────────────────────────
    cv_score, cv_sig = _cv_section_structure(text)

    # ── 6. Vague achievement claims ───────────────────────────────────────
    ach_score, ach_sig = _achievement_claim_density(text)

    # ── 7. Skill list density ─────────────────────────────────────────────
    skill_score, skill_sig = _skill_list_density(text)

    # ── 8. English-specific linguistic heuristics ─────────────────────────
    phrase_score, phrase_sig = _transition_phrases(text)
    action_score, action_sig = _action_verb_density(text)
    cover_score,  cover_sig  = _cover_letter_patterns(text)

    signals += [
        perp_sig, tmpl_sig, buzz_sig, cv_sig, ach_sig, skill_sig,
        burst_sig, vocab_sig, phrase_sig, para_sig, sim_sig,
        action_sig, cover_sig, section_sig,
    ]

    # Perplexity: floor at 0.38 — jargon/structured text inflates GPT-2 perplexity
    p = max(perp_score, 0.38) if perp_score is not None else 0.45

    h_scores = [
        tmpl_score, buzz_score, cv_score, ach_score, skill_score,
        burst_score, vocab_score, phrase_score, para_score, sim_score,
        action_score, cover_score,
    ]

    # Count how many signals strongly indicate AI content
    high_count = sum(1 for s in h_scores if s >= 0.65)

    if ml_prob is not None and ml_prob >= 0.25:
        # RoBERTa has a meaningful signal — use full ML blend
        # RoBERTa 15% + perplexity 20% + heuristics 65%
        h_weights = [0.10, 0.08, 0.07, 0.06, 0.05, 0.06, 0.04, 0.06, 0.02, 0.02, 0.05, 0.04]  # sum=0.65
        h_prob = sum(s * w for s, w in zip(h_scores, h_weights))
        ai_prob = ml_prob * 0.15 + p * 0.20 + h_prob
    else:
        # RoBERTa unavailable or unreliable on this text type
        h_weights = [0.15, 0.11, 0.10, 0.08, 0.08, 0.09, 0.07, 0.08, 0.03, 0.03, 0.07, 0.06]  # sum=0.95
        h_prob = sum(s * w for s, w in zip(h_scores, h_weights))
        ai_prob = h_prob * 0.80 + p * 0.20

    # Section uniformity adds an independent boost for English docs too
    if section_score >= 0.72:
        ai_prob = max(ai_prob, min(ai_prob + (section_score - 0.72) * 0.5, 1.0))

    # Signal consensus floor
    if high_count >= 6:
        ai_prob = max(ai_prob, 0.88)
    elif high_count >= 5:
        ai_prob = max(ai_prob, 0.82)
    elif high_count >= 4:
        ai_prob = max(ai_prob, 0.75)
    elif high_count >= 3:
        ai_prob = max(ai_prob, 0.65)

    # Individual strong signal floors (independent — not elif so multiple can fire)
    if tmpl_score >= 0.80:
        ai_prob = max(ai_prob, 0.78)
    if cover_score >= 0.70 and (buzz_score >= 0.50 or phrase_score >= 0.45):
        ai_prob = max(ai_prob, 0.72)
    if cv_score >= 0.75 and buzz_score >= 0.55:
        ai_prob = max(ai_prob, 0.72)
    if buzz_score >= 0.80 and phrase_score >= 0.45:
        ai_prob = max(ai_prob, 0.68)
    if ach_score >= 0.75 and skill_score >= 0.55:
        ai_prob = max(ai_prob, 0.65)
    if action_score >= 0.75 and (cv_score >= 0.50 or buzz_score >= 0.50):
        ai_prob = max(ai_prob, 0.68)

    return float(min(max(ai_prob, 0.0), 1.0)), signals


# ─────────────────────────────────────────────────────────────────────────────
# AUDIO HELPERS
# ─────────────────────────────────────────────────────────────────────────────

# Known AI audio/music generation tool name fragments (lowercase bytes for binary scan)
_AI_AUDIO_TOOL_MARKERS = [
    # Music generation platforms
    b"suno", b"udio", b"stable audio", b"musicgen", b"audiogen",
    b"mubert", b"aiva", b"soundraw", b"boomy", b"loudly",
    b"beatoven", b"riffusion", b"musiclm", b"jukebox", b"musicfy",
    b"soundful", b"ecrett", b"amper",
    # Voice / TTS
    b"elevenlabs", b"eleven labs", b"bark tts", b"vall-e",
    b"tortoisetts", b"coqui", b"playht", b"play.ht",
    b"resemble.ai", b"descript", b"speechify", b"voicify",
    # Generic AI tag markers
    b"generated by ai", b"ai generated", b"ai-generated",
    b"created by ai", b"ai music", b"created with ai",
]


def _audio_ai_tool_scan(file_path: Path) -> Tuple[float, DetectionSignal]:
    """
    Binary scan of the file + full metadata tag scan for known AI
    music / TTS generation tool markers.

    This is the highest-confidence signal: a match is near-definitive.
    """
    found_binary: List[str] = []
    found_tags: List[str] = []

    # ── Binary scan (first 64 KB + last 8 KB) ─────────────────────────────
    try:
        size = file_path.stat().st_size
        with open(file_path, "rb") as f:
            head = f.read(min(65536, size))
            tail = b""
            if size > 65536:
                f.seek(max(0, size - 8192))
                tail = f.read(8192)
        blob = (head + tail).lower()
        for marker in _AI_AUDIO_TOOL_MARKERS:
            if marker in blob:
                found_binary.append(marker.decode("utf-8", errors="ignore"))
    except Exception:
        pass

    # ── Full metadata tag scan ─────────────────────────────────────────────
    try:
        from mutagen import File as MutagenFile
        from app.services.metadata import _flag_software
        audio = MutagenFile(file_path)
        if audio and audio.tags:
            tag_text = " ".join(str(v) for v in audio.tags.values()).lower()
            for marker in _AI_AUDIO_TOOL_MARKERS:
                m = marker.decode("utf-8", errors="ignore")
                if m in tag_text and m not in found_tags:
                    found_tags.append(m)
            # Also run the generic software-flag check on encoder fields
            for k in ("TSSE", "TENC", "encoder", "encoded_by", "comment", "COMM"):
                if k in audio.tags:
                    val = str(audio.tags[k])
                    susp, _note = _flag_software(val)
                    if susp:
                        entry = f"{k}={val[:60]}"
                        if entry not in found_tags:
                            found_tags.append(entry)
    except Exception:
        pass

    all_found = list(dict.fromkeys(found_binary + found_tags))  # deduplicate, keep order

    if all_found:
        score, severity = 0.95, SignalSeverity.HIGH
        desc = (
            f"Known AI audio generation tool detected: {', '.join(all_found[:4])}. "
            "This is a near-definitive indicator of AI-generated audio."
        )
    else:
        score, severity = 0.20, SignalSeverity.LOW
        desc = "No known AI audio generation tool markers found in file binary or metadata."

    return score, DetectionSignal(
        name="AI Audio Tool Detection",
        description=desc,
        severity=severity,
        score=score,
        details=f"Binary hits: {found_binary or 'none'} | Tag hits: {found_tags or 'none'}",
    )


def _audio_metadata_check(file_path: Path) -> Tuple[float, DetectionSignal]:
    """
    Checks metadata *completeness* for a music file.

    AI music platforms typically export tracks with missing or minimal tags
    (no artist, no album, no genre, no year).  A real studio track always
    has at least artist + title.  Missing 3+ of the 5 key fields is a strong
    signal.  Also scans comment/description text for AI-related language.
    """
    try:
        from mutagen import File as MutagenFile
        audio = MutagenFile(file_path)
        if audio is None:
            return 0.55, DetectionSignal(
                name="Audio Metadata",
                description="Could not parse audio file metadata.",
                severity=SignalSeverity.MEDIUM,
                score=0.55,
            )

        if not audio.tags:
            return 0.68, DetectionSignal(
                name="Audio Metadata Absent",
                description=(
                    "No metadata tags found. Professional music always carries at least "
                    "artist and title tags; AI music tools frequently skip tagging."
                ),
                severity=SignalSeverity.HIGH,
                score=0.68,
            )

        tags = audio.tags

        def _has_field(*keys: str) -> bool:
            for k in keys:
                for tag_key in tags.keys():
                    if k.lower() in str(tag_key).lower():
                        val = str(tags[tag_key]).strip()
                        if val and val not in ("", "0", "unknown", "none"):
                            return True
            return False

        has_artist = _has_field("artist", "TPE1", "TPE2", "author")
        has_title  = _has_field("title",  "TIT2", "TIT1")
        has_album  = _has_field("album",  "TALB")
        has_year   = _has_field("year",   "date", "TDRC", "TYER")
        has_genre  = _has_field("genre",  "TCON")

        missing = sum([not has_artist, not has_title, not has_album, not has_year, not has_genre])

        # Scan all tag text for AI-related language
        tag_text = " ".join(str(v) for v in tags.values()).lower()
        ai_lang  = ["generated", "ai ", "artificial", "synthetic", "suno", "udio",
                    "prompt:", "style:", "created with", "instrumental"]
        comment_suspicious = any(p in tag_text for p in ai_lang)

        field_str = (
            f"artist={'✓' if has_artist else '✗'} "
            f"title={'✓' if has_title else '✗'} "
            f"album={'✓' if has_album else '✗'} "
            f"year={'✓' if has_year else '✗'} "
            f"genre={'✓' if has_genre else '✗'}"
        )

        if missing >= 4 or (missing >= 3 and comment_suspicious):
            score, severity = 0.75, SignalSeverity.HIGH
            desc = (
                f"Severely incomplete metadata ({missing}/5 key fields missing). "
                f"{field_str}. AI-generated audio rarely carries full music metadata."
            )
        elif missing >= 3:
            score, severity = 0.58, SignalSeverity.MEDIUM
            desc = f"Incomplete metadata ({missing}/5 fields missing). {field_str}."
        elif missing >= 2 or comment_suspicious:
            score, severity = 0.42, SignalSeverity.MEDIUM
            if comment_suspicious:
                desc = f"Metadata contains AI-related language. {field_str}."
            else:
                desc = f"Partial metadata ({missing}/5 fields missing). {field_str}."
        else:
            score, severity = 0.12, SignalSeverity.LOW
            desc = f"Complete music metadata present. {field_str}."

        return score, DetectionSignal(
            name="Audio Metadata Completeness",
            description=desc,
            severity=severity,
            score=score,
            details=field_str,
        )
    except Exception as exc:
        return 0.40, DetectionSignal(
            name="Audio Metadata Error",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.40,
        )


def _audio_format_check(file_path: Path) -> Tuple[float, DetectionSignal]:
    """
    Format heuristics: sample rate, duration, bitrate, channel count.

    Key insight: TTS tools use 16/22/24 kHz; AI *music* tools (Suno, Udio)
    export at 44100/48000 Hz — so sample rate alone is insufficient.
    We therefore score a combination of factors.
    """
    try:
        from mutagen import File as MutagenFile
        audio = MutagenFile(file_path)
        if audio is None or not hasattr(audio, "info"):
            return 0.40, DetectionSignal(
                name="Audio Format",
                description="Cannot determine audio format.",
                severity=SignalSeverity.LOW,
                score=0.40,
            )

        info     = audio.info
        duration = getattr(info, "length",      0)
        bitrate  = getattr(info, "bitrate",     0)
        sr       = getattr(info, "sample_rate", 0)
        channels = getattr(info, "channels",    0)

        suspicion = 0.0
        flags: List[str] = []

        # TTS-typical sample rates
        if sr in (16000, 22050, 24000):
            suspicion += 0.35
            flags.append(f"TTS sample rate ({sr} Hz)")

        # Very short — TTS demo or AI preview
        if 0 < duration < 15:
            suspicion += 0.22
            flags.append(f"Short duration ({duration:.1f}s)")

        # AI generation platforms often cap at exact multiples of 30s / 60s
        if duration > 0 and abs(duration - round(duration)) < 0.08 and round(duration) % 30 == 0:
            suspicion += 0.15
            flags.append(f"Exact duration ({duration:.0f}s — common AI generation limit)")

        # Mono + TTS sample rate
        if channels == 1 and sr in (16000, 22050, 24000):
            suspicion += 0.18
            flags.append("Mono + TTS sample rate")

        score = min(suspicion, 0.85) if flags else 0.15

        if score >= 0.55:
            severity = SignalSeverity.HIGH
            desc = f"Audio format strongly consistent with AI synthesis: {'; '.join(flags)}."
        elif score >= 0.28:
            severity = SignalSeverity.MEDIUM
            desc = f"Suspicious format properties: {'; '.join(flags)}."
        else:
            severity = SignalSeverity.LOW
            desc = (
                f"Normal audio format (SR={sr} Hz, {duration:.1f}s, "
                f"{bitrate // 1000 if bitrate else '?'}kbps). No AI format markers."
            )

        return score, DetectionSignal(
            name="Audio Format Properties",
            description=desc,
            severity=severity,
            score=score,
            details=f"Duration={duration:.1f}s | SR={sr} Hz | Bitrate={bitrate} bps | Channels={channels}",
        )
    except Exception as exc:
        return 0.40, DetectionSignal(
            name="Audio Format Error",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.40,
        )


def _decode_audio_float32(file_path: Path):
    """
    Decode any audio file to float32 mono numpy array.

    Strategy:
      1. Try soundfile (fast; handles WAV, FLAC, OGG but NOT MP3/AAC).
      2. Fall back to ffmpeg subprocess for compressed formats.

    Returns (data: np.ndarray | None, sample_rate: int).
    """
    import numpy as np

    # ── soundfile (fast path) ──────────────────────────────────────────────
    try:
        import soundfile as sf
        data, sr = sf.read(str(file_path), always_2d=False)
        if data.ndim > 1:
            data = data.mean(axis=1)
        return data.astype(np.float32), int(sr)
    except Exception:
        pass

    # ── ffmpeg subprocess (handles MP3, AAC, M4A, …) ──────────────────────
    try:
        import subprocess
        result = subprocess.run(
            [
                "ffmpeg", "-hide_banner", "-loglevel", "error",
                "-i", str(file_path),
                "-ac", "1",        # mono
                "-ar", "22050",    # 22 kHz — sufficient for our analysis
                "-f", "f32le",     # raw float32 little-endian PCM
                "pipe:1",
            ],
            capture_output=True,
            timeout=45,
        )
        if result.returncode == 0 and result.stdout:
            data = np.frombuffer(result.stdout, dtype=np.float32).copy()
            return data, 22050
    except Exception:
        pass

    return None, 0


def _audio_spectral_analysis(file_path: Path) -> Tuple[float, DetectionSignal]:
    """
    Frame-level spectral consistency analysis — now works on MP3/AAC via ffmpeg.

    AI-generated audio (both TTS and AI music) tends to exhibit:
      - Unnaturally low variance in zero-crossing rate across frames
      - Very stable short-term energy envelope
      - Abnormally consistent spectral centroid

    Natural recordings vary considerably in all three dimensions.
    """
    import numpy as np

    try:
        data, sr = _decode_audio_float32(file_path)
        if data is None or len(data) == 0:
            return 0.45, DetectionSignal(
                name="Audio Spectral Analysis",
                description="Could not decode audio for spectral analysis.",
                severity=SignalSeverity.LOW,
                score=0.45,
            )

        max_val = np.abs(data).max()
        if max_val > 0:
            data /= max_val

        frame_size = max(int(sr * 0.025), 64)   # 25 ms
        hop_size   = max(int(sr * 0.010), 32)    # 10 ms

        if len(data) < frame_size * 10:
            return 0.40, DetectionSignal(
                name="Audio Spectral Analysis",
                description="Audio too short for spectral analysis.",
                severity=SignalSeverity.LOW,
                score=0.40,
            )

        # Analyse up to 30 s worth of frames for efficiency
        max_frames  = int(30 * sr / hop_size)
        window      = np.hanning(frame_size)
        freqs       = np.fft.rfftfreq(frame_size, d=1.0 / sr)
        zcr_vals, energy_vals, centroid_vals = [], [], []

        frame_count = 0
        for i in range(0, len(data) - frame_size, hop_size):
            if frame_count >= max_frames:
                break
            frame = data[i: i + frame_size]
            zcr_vals.append(float(np.mean(np.diff(np.sign(frame)) != 0)))
            energy_vals.append(float(np.sqrt(np.mean(frame ** 2))))
            mag   = np.abs(np.fft.rfft(frame * window))
            total = mag.sum()
            centroid_vals.append(float((freqs * mag).sum() / (total + 1e-9)))
            frame_count += 1

        def _cv(vals: list) -> float:
            a = np.array(vals)
            return float(a.std() / (a.mean() + 1e-9))

        zcr_cv      = _cv(zcr_vals)
        energy_cv   = _cv(energy_vals)
        centroid_cv = _cv(centroid_vals)
        consistency = (zcr_cv + energy_cv + centroid_cv) / 3.0

        # Tighter thresholds — AI music is more dynamic than TTS but still
        # more homogeneous than organic recordings
        if consistency < 0.30:
            score, severity = 0.85, SignalSeverity.HIGH
            desc = (
                f"Highly consistent spectral profile (composite CV={consistency:.3f}: "
                f"ZCR={zcr_cv:.2f}, Energy={energy_cv:.2f}, Centroid={centroid_cv:.2f}). "
                "AI synthesis produces unnaturally uniform audio dynamics."
            )
        elif consistency < 0.55:
            score, severity = 0.58, SignalSeverity.MEDIUM
            desc = (
                f"Moderately consistent audio profile (composite CV={consistency:.3f}: "
                f"ZCR={zcr_cv:.2f}, Energy={energy_cv:.2f}, Centroid={centroid_cv:.2f})."
            )
        elif consistency < 0.80:
            score, severity = 0.30, SignalSeverity.LOW
            desc = (
                f"Some consistency detected (composite CV={consistency:.3f}) but within "
                "normal range for produced/mastered music."
            )
        else:
            score, severity = 0.12, SignalSeverity.LOW
            desc = (
                f"Natural audio variability (composite CV={consistency:.3f}). "
                "Spectral profile shifts naturally — consistent with organic recording."
            )

        return score, DetectionSignal(
            name="Audio Spectral Analysis",
            description=desc,
            severity=severity,
            score=score,
            details=(
                f"ZCR CV={zcr_cv:.3f} | Energy CV={energy_cv:.3f} | "
                f"Centroid CV={centroid_cv:.3f} | Composite={consistency:.3f} | "
                f"Duration={len(data) / sr:.1f}s | SR={sr} Hz | Frames={frame_count}"
            ),
        )

    except Exception as exc:
        return 0.40, DetectionSignal(
            name="Audio Spectral Analysis",
            description=f"Spectral analysis failed: {exc}",
            severity=SignalSeverity.LOW,
            score=0.40,
        )


def _audio_dynamic_range(file_path: Path) -> Tuple[float, DetectionSignal]:
    """
    Dynamic range analysis — AI music is characteristically over-compressed.

    AI music generation platforms apply aggressive loudness normalisation
    (LUFS targeting), resulting in very low dynamic range and unusually
    consistent RMS levels compared to organic studio recordings (which
    typically span 15–30+ dB of dynamic headroom).
    """
    import numpy as np

    try:
        data, sr = _decode_audio_float32(file_path)
        if data is None or len(data) < sr * 2:
            return 0.40, DetectionSignal(
                name="Dynamic Range Analysis",
                description="Insufficient audio for dynamic range analysis.",
                severity=SignalSeverity.LOW,
                score=0.40,
            )

        # RMS in 500 ms windows with 50 % overlap
        win = max(int(sr * 0.5), 1)
        rms_vals = []
        for i in range(0, len(data) - win, win // 2):
            chunk = data[i: i + win]
            rms = float(np.sqrt(np.mean(chunk ** 2)))
            if rms > 0.001:   # skip near-silence
                rms_vals.append(rms)

        if len(rms_vals) < 4:
            return 0.40, DetectionSignal(
                name="Dynamic Range Analysis",
                description="Audio too short or silent for dynamic range analysis.",
                severity=SignalSeverity.LOW,
                score=0.40,
            )

        rms_arr          = np.array(rms_vals)
        dynamic_range_db = 20 * np.log10(rms_arr.max() / (rms_arr.min() + 1e-9))
        rms_cv           = float(rms_arr.std() / (rms_arr.mean() + 1e-9))

        # Flat / near-silence sections — padding common in AI generation
        total_wins = max(len(data) // (win // 2), 1)
        flat_ratio = (total_wins - len(rms_vals)) / total_wins

        suspicion = 0.0
        flags: List[str] = []

        if dynamic_range_db < 8:
            suspicion += 0.42
            flags.append(f"Very low dynamic range ({dynamic_range_db:.1f} dB)")
        elif dynamic_range_db < 14:
            suspicion += 0.22
            flags.append(f"Low dynamic range ({dynamic_range_db:.1f} dB)")

        if rms_cv < 0.15:
            suspicion += 0.38
            flags.append(f"Hyper-consistent loudness (CV={rms_cv:.3f})")
        elif rms_cv < 0.30:
            suspicion += 0.18
            flags.append(f"Low loudness variation (CV={rms_cv:.3f})")

        if flat_ratio > 0.15:
            suspicion += 0.18
            flags.append(f"High silent-frame ratio ({flat_ratio:.0%})")

        score = min(suspicion, 0.88)

        if score >= 0.55:
            severity = SignalSeverity.HIGH
            desc = f"Hyper-compressed dynamics — hallmark of AI music generation: {'; '.join(flags)}."
        elif score >= 0.28:
            severity = SignalSeverity.MEDIUM
            desc = f"Suspicious dynamic compression: {'; '.join(flags)}." if flags else "Moderately compressed dynamics."
        else:
            severity = SignalSeverity.LOW
            desc = (
                f"Natural dynamic range ({dynamic_range_db:.1f} dB, RMS CV={rms_cv:.3f}). "
                "No unusual compression detected."
            )

        return score, DetectionSignal(
            name="Dynamic Range Analysis",
            description=desc,
            severity=severity,
            score=score,
            details=f"DR={dynamic_range_db:.1f}dB | RMS CV={rms_cv:.3f} | Flat ratio={flat_ratio:.1%}",
        )

    except Exception as exc:
        return 0.35, DetectionSignal(
            name="Dynamic Range Analysis",
            description=f"Dynamic range analysis failed: {exc}",
            severity=SignalSeverity.LOW,
            score=0.35,
        )


def _audio_30s_deep_scan(file_path: Path) -> Tuple[float, DetectionSignal]:
    """
    Deep acoustic analysis over a 30-second segment.

    Extracts 30 s starting at 5 s in (skipping intros/silence) and computes:

    1. Spectral flatness (Wiener entropy) CV — AI synthesis has unnaturally
       stable flatness across frames; natural recordings vary widely.
    2. Sub-band energy CV (bass/mid/high) — AI music maintains suspiciously
       consistent energy distribution across frequency bands over time.
    3. Amplitude envelope CV — AI music loudness barely fluctuates at the
       frame level; real recordings breathe with natural micro-dynamics.

    All three metrics are measured as coefficient of variation (CV = std/mean).
    Low CV across frames = too consistent = AI-like.
    """
    import numpy as np
    import subprocess

    def _extract_segment(path: Path, skip: float, duration: float):
        """Decode `duration` seconds starting at `skip` to float32 mono @ 22050 Hz."""
        cmd = [
            "ffmpeg", "-hide_banner", "-loglevel", "error",
            "-ss", str(skip),
            "-i", str(path),
            "-t", str(duration),
            "-ac", "1",
            "-ar", "22050",
            "-f", "f32le",
            "pipe:1",
        ]
        r = subprocess.run(cmd, capture_output=True, timeout=60)
        if r.returncode == 0 and r.stdout:
            return np.frombuffer(r.stdout, dtype=np.float32).copy()
        return None

    try:
        sr = 22050
        # Try skipping the first 5 s; fall back to start of file
        data = _extract_segment(file_path, skip=5.0, duration=30.0)
        if data is None or len(data) < sr * 5:
            data = _extract_segment(file_path, skip=0.0, duration=30.0)
        if data is None or len(data) < sr * 2:
            return 0.45, DetectionSignal(
                name="30s Deep Acoustic Scan",
                description="Could not extract audio segment for deep analysis.",
                severity=SignalSeverity.LOW,
                score=0.45,
            )

        max_val = np.abs(data).max()
        if max_val > 0:
            data /= max_val

        frame_size = 2048          # ~93 ms @ 22050 Hz — good frequency resolution
        hop_size   = 512           # ~23 ms hop
        window     = np.hanning(frame_size)
        freqs      = np.fft.rfftfreq(frame_size, d=1.0 / sr)

        bass_mask = freqs < 300
        mid_mask  = (freqs >= 300)  & (freqs < 3000)
        high_mask = freqs >= 3000

        flatness_vals: list  = []
        envelope_vals: list  = []
        bass_e: list         = []
        mid_e: list          = []
        high_e: list         = []

        for i in range(0, len(data) - frame_size, hop_size):
            if len(flatness_vals) >= 2000:
                break
            frame = data[i: i + frame_size] * window

            # Amplitude envelope
            envelope_vals.append(float(np.sqrt(np.mean(frame ** 2))))

            mag = np.abs(np.fft.rfft(frame)) + 1e-10

            # Wiener spectral flatness
            log_mean   = float(np.mean(np.log(mag)))
            arith_mean = float(np.mean(mag))
            flatness_vals.append(float(np.exp(log_mean) / (arith_mean + 1e-10)))

            # Sub-band energies
            bass_e.append(float(mag[bass_mask].mean()))
            mid_e.append(float(mag[mid_mask].mean()))
            high_e.append(float(mag[high_mask].mean()))

        def _cv(vals: list) -> float:
            a = np.array(vals)
            m = float(a.mean())
            return float(a.std() / (m + 1e-9)) if m > 0 else 1.0

        envelope_cv   = _cv(envelope_vals)
        flatness_cv   = _cv(flatness_vals)
        bass_cv       = _cv(bass_e)
        mid_cv        = _cv(mid_e)
        high_cv       = _cv(high_e)
        band_cv       = (bass_cv + mid_cv + high_cv) / 3.0
        mean_flatness = float(np.mean(flatness_vals))

        suspicion = 0.0
        flags: List[str] = []

        # Envelope consistency
        if envelope_cv < 0.18:
            suspicion += 0.36
            flags.append(f"hyper-stable loudness envelope (CV={envelope_cv:.3f})")
        elif envelope_cv < 0.35:
            suspicion += 0.16
            flags.append(f"stable loudness envelope (CV={envelope_cv:.3f})")

        # Spectral flatness consistency
        if flatness_cv < 0.08:
            suspicion += 0.32
            flags.append(f"ultra-consistent spectral flatness (CV={flatness_cv:.3f})")
        elif flatness_cv < 0.18:
            suspicion += 0.16
            flags.append(f"consistent spectral flatness (CV={flatness_cv:.3f})")

        # Sub-band energy consistency
        if band_cv < 0.22:
            suspicion += 0.28
            flags.append(f"rigid sub-band energy (CV={band_cv:.3f})")
        elif band_cv < 0.38:
            suspicion += 0.12
            flags.append(f"stable sub-band energy (CV={band_cv:.3f})")

        score = min(suspicion, 0.92)

        if score >= 0.60:
            severity = SignalSeverity.HIGH
            desc = f"30s deep scan: strong AI acoustic signatures — {'; '.join(flags)}."
        elif score >= 0.32:
            severity = SignalSeverity.MEDIUM
            desc = f"30s deep scan: suspicious acoustic patterns — {'; '.join(flags)}."
        else:
            severity = SignalSeverity.LOW
            desc = (
                f"30s deep scan: natural acoustic variation detected "
                f"(envelope CV={envelope_cv:.3f}, flatness CV={flatness_cv:.3f}, "
                f"band CV={band_cv:.3f})."
            )

        return score, DetectionSignal(
            name="30s Deep Acoustic Scan",
            description=desc,
            severity=severity,
            score=score,
            details=(
                f"Envelope CV={envelope_cv:.3f} | Flatness CV={flatness_cv:.3f} | "
                f"Band CV={band_cv:.3f} (bass={bass_cv:.3f} mid={mid_cv:.3f} high={high_cv:.3f}) | "
                f"Mean flatness={mean_flatness:.4f} | Frames={len(flatness_vals)} | SR={sr} Hz"
            ),
        )

    except Exception as exc:
        return 0.40, DetectionSignal(
            name="30s Deep Acoustic Scan",
            description=f"Deep scan failed: {exc}",
            severity=SignalSeverity.LOW,
            score=0.40,
        )


# ─────────────────────────────────────────────────────────────────────────────
# VIDEO HELPERS
# ─────────────────────────────────────────────────────────────────────────────

async def _extract_and_analyze_frame(file_path: Path) -> Tuple[float, List[DetectionSignal]]:
    """Extract first readable JPEG frame from video and run image analysis."""
    try:
        import numpy as np
        from PIL import Image

        with open(file_path, "rb") as f:
            data = f.read(1024 * 1024)

        jpg_start = data.find(b"\xff\xd8\xff")
        if jpg_start != -1:
            jpg_end = data.find(b"\xff\xd9", jpg_start)
            if jpg_end != -1 and jpg_end - jpg_start > 1000:
                frame_data = data[jpg_start:jpg_end + 2]
                img = Image.open(io.BytesIO(frame_data))
                arr = np.array(img.convert("RGB"), dtype=np.float32)
                s_ela,   sig_ela   = _ela_analysis(img)
                s_noise, sig_noise = _noise_analysis(arr)
                s_freq,  sig_freq  = _frequency_analysis(arr)
                prob = s_ela * 0.40 + s_noise * 0.30 + s_freq * 0.30
                return prob, [sig_ela, sig_noise, sig_freq]
    except Exception:
        pass

    return 0.35, [DetectionSignal(
        name="Video Frame Analysis",
        description="Could not extract frames for deep analysis. Metadata-only assessment performed.",
        severity=SignalSeverity.LOW,
        score=0.35,
    )]


def _video_container_check(file_path: Path) -> Tuple[float, DetectionSignal]:
    try:
        with open(file_path, "rb") as f:
            header = f.read(512)

        AI_VIDEO_MARKERS = [b"RunwayML", b"Pika", b"Sora", b"Kling", b"Stable Video"]
        for marker in AI_VIDEO_MARKERS:
            if marker.lower() in header.lower():
                return 0.90, DetectionSignal(
                    name="AI Video Tool Detected",
                    description=f"Container header references known AI video tool: {marker.decode()}",
                    severity=SignalSeverity.HIGH,
                    score=0.90,
                )

        return 0.18, DetectionSignal(
            name="Video Container",
            description="No AI video tool markers found in container header.",
            severity=SignalSeverity.LOW,
            score=0.18,
        )
    except Exception as exc:
        return 0.30, DetectionSignal(
            name="Video Container Check Failed",
            description=str(exc),
            severity=SignalSeverity.LOW,
            score=0.30,
        )
